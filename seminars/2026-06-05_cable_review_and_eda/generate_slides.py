"""Generate English PPT for 2026-06-05 group seminar (9 slides).

Slides:
    1. Cover · 2. Field overview · 3. Zoom-in · 4. Feasible path
    5. Beier 2023 · 6. Beyond Beier
    7. AI purpose + data + flowchart
    8. EDA findings (Plot 1 + 2)
    9. EDA findings (Plot 3 + 4)

Paths are relative to this script — run `python generate_slides.py` from
anywhere; PNGs are pulled from `<repo>/figures/eda_wind_scada/` and
`slides_en.pptx` is written next to this script.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent          # seminars/<date>_<topic>/
ROOT = HERE.parent.parent                       # repo root
PNG_DIR = ROOT / "figures" / "eda_wind_scada"
SLIDES_OUT = HERE / "slides_en.pptx"

FONT = "Calibri"
COLOR_TITLE = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
COLOR_HEADER_BG = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ROW_ALT = RGBColor(0xF2, 0xF2, 0xF2)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOFT = RGBColor(0x66, 0x66, 0x66)
COLOR_HIGHLIGHT_BG = RGBColor(0xFF, 0xF5, 0xE6)
COLOR_HIGHLIGHT_BORDER = RGBColor(0xC0, 0x39, 0x2B)
COLOR_COVER_BG = RGBColor(0x1F, 0x3A, 0x5F)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, text, size=12, bold=False, color=COLOR_TEXT, italic=False):
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=COLOR_TEXT, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color, italic=italic)
    return tb


def add_bullets(slide, left, top, width, height, items, size=13,
                bullet="• ", line_color=COLOR_TEXT, space_after=4):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        if isinstance(item, str):
            r = p.add_run()
            set_run(r, bullet + item, size=size, color=line_color)
        else:
            prefix, bold_part, rest = item
            r0 = p.add_run()
            set_run(r0, bullet + prefix, size=size, color=line_color)
            r1 = p.add_run()
            set_run(r1, bold_part, size=size, bold=True, color=COLOR_ACCENT)
            r2 = p.add_run()
            set_run(r2, rest, size=size, color=line_color)
    return tb


def add_table(slide, left, top, width, height, data, col_widths,
              header_size=13, body_size=13):
    rows = len(data)
    cols = len(data[0])
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w
    for ri in range(rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            cell.margin_left = Emu(50000)
            cell.margin_right = Emu(50000)
            cell.margin_top = Emu(40000)
            cell.margin_bottom = Emu(40000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_HEADER_BG
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            content = data[ri][ci]
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ri == 0 else PP_ALIGN.LEFT
            if isinstance(content, str):
                content = [(content, False)]
            for txt, b in content:
                r = p.add_run()
                if ri == 0:
                    set_run(r, txt, size=header_size, bold=True,
                            color=COLOR_HEADER_FG)
                else:
                    col = COLOR_ACCENT if b else COLOR_TEXT
                    set_run(r, txt, size=body_size, bold=b, color=col)
    return shape


def add_highlight(slide, left, top, width, height, title, body,
                  title_size=14, body_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_HIGHLIGHT_BG
    shape.line.color.rgb = COLOR_HIGHLIGHT_BORDER
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(100000)
    tf.margin_right = Emu(100000)
    tf.margin_top = Emu(70000)
    tf.margin_bottom = Emu(70000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r, title, size=title_size, bold=True, color=COLOR_ACCENT)
    for line in body:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        set_run(r2, line, size=body_size, color=COLOR_TEXT)
    return shape


def add_footer(slide, idx, total=9, section="Ocean Engineering"):
    add_textbox(slide, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.25),
                f"P{idx} / {total}  ·  {section}",
                size=9, color=COLOR_SOFT, align=PP_ALIGN.RIGHT)


# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ============================================================
# Cover
# ============================================================
cover = prs.slides.add_slide(blank)

# Background band on left
band = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(0),
                              Inches(4.0), SLIDE_H)
band.fill.solid()
band.fill.fore_color.rgb = COLOR_COVER_BG
band.line.fill.background()
band.shadow.inherit = False

# Vertical accent line
accent = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(4.0), Inches(0),
                                Inches(0.08), SLIDE_H)
accent.fill.solid()
accent.fill.fore_color.rgb = COLOR_ACCENT
accent.line.fill.background()
accent.shadow.inherit = False

# Date / event tag (on left band, vertical)
add_textbox(cover, Inches(0.5), Inches(0.6), Inches(3.0), Inches(0.4),
            "SEMINAR", size=14, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF))
add_textbox(cover, Inches(0.5), Inches(1.05), Inches(3.0), Inches(0.4),
            "20260605", size=20, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF))

# Presenter (bottom of left band)
add_textbox(cover, Inches(0.5), Inches(6.2), Inches(3.0), Inches(0.4),
            "Presenter", size=11,
            color=RGBColor(0xCC, 0xCC, 0xCC))
add_textbox(cover, Inches(0.5), Inches(6.55), Inches(3.0), Inches(0.5),
            "Diao Jingyu", size=20, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF))

# Main title (right side)
add_textbox(cover, Inches(4.7), Inches(2.4), Inches(8.3), Inches(0.5),
            "Ocean Engineering Section",
            size=14, bold=True, color=COLOR_ACCENT)
add_textbox(cover, Inches(4.7), Inches(2.85), Inches(8.3), Inches(1.4),
            "Dynamic Power Cables for\nFloating Offshore Wind Turbines",
            size=34, bold=True, color=COLOR_TITLE)
add_textbox(cover, Inches(4.7), Inches(4.6), Inches(8.3), Inches(0.45),
            "From literature review to a research direction",
            size=16, italic=True, color=COLOR_SOFT)

# Roadmap chips at bottom right
add_textbox(cover, Inches(4.7), Inches(5.55), Inches(8.3), Inches(0.3),
            "Roadmap",
            size=11, bold=True, color=COLOR_SOFT)
add_textbox(cover, Inches(4.7), Inches(5.85), Inches(8.3), Inches(0.9),
            "1. Field overview   →   2. Accessible directions   →   "
            "3. Our feasible path   →   4. Beier 2023   →   5. Beyond Beier",
            size=11, color=COLOR_TEXT)


# ============================================================
# P1: Field Overview — six research areas from the review
# ============================================================
s1 = prs.slides.add_slide(blank)

add_textbox(s1, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.55),
            "Field Overview — Six Research Areas",
            size=28, bold=True, color=COLOR_TITLE)
add_textbox(s1, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4),
            "Cerik & Huang, Ocean Engineering 311 (2024) — a review of "
            "the last ~20 years",
            size=13, italic=True, color=COLOR_SOFT)

overview_header = ["#", "Research Area", "Key Question", "State of the Art",
                   "Accessible to us?"]
overview = [
    ["1",
     [("Mechanical properties", True)],
     "Stiffness, stick-slip nonlinearity?",
     "Bending hysteresis quantified (Hu 2022; Ménard 2023)",
     "❌ needs bending rig"],
    ["2",
     [("Failure mechanisms", True)],
     "Why does it fail?",
     "Fretting, water tree, biofouling — identified separately",
     "❌ needs FEA / electrical lab"],
    ["3",
     [("Fatigue analysis", True)],
     "How long can it last?",
     "4-step workflow standardised (global → local → rainflow → Miner)",
     [("✅ global half", True)]],
    ["4",
     [("Experimental studies", True)],
     "Where do test data come from?",
     "Vendors hold the data; public datasets are scarce",
     "❌ needs full-scale rig"],
    ["5",
     [("Local (cross-section) analysis", True)],
     "Internal stress distribution?",
     "RVE / beam / homogenisation strategies mature",
     "❌ needs UFLEX / Helica"],
    ["6",
     [("Global load analysis", True)],
     "Overall response of the whole cable?",
     "OrcaFlex dominates; OpenFAST + MoorDyn now open-source",
     [("✅ our core tool", True)]],
]
add_table(s1, Inches(0.4), Inches(1.45),
          Inches(12.5), Inches(4.5),
          [overview_header] + overview,
          [Inches(0.4), Inches(2.4), Inches(3.0), Inches(4.6), Inches(2.1)],
          header_size=13, body_size=13)

add_highlight(s1, Inches(0.4), Inches(6.15),
              Inches(12.5), Inches(0.85),
              "Filter result → next slide",
              [
                  "Only areas 3 (global half) and 6 are accessible. "
                  "On P3 we zoom into these two areas and split them into 4 active sub-directions.",
              ],
              title_size=14, body_size=13)

add_footer(s1, 2)


# ============================================================
# P2: Accessible Directions (was P1, simplified)
# ============================================================
s2 = prs.slides.add_slide(blank)

add_textbox(s2, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.55),
            "Zoom in — 4 Sub-directions Within Accessible Areas",
            size=28, bold=True, color=COLOR_TITLE)
add_textbox(s2, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4),
            "Expanding P2's Area 3 (Fatigue, global half) + Area 6 (Global analysis)",
            size=13, italic=True, color=COLOR_SOFT)

# Filter breadcrumb
add_highlight(s2, Inches(0.4), Inches(1.3),
              Inches(12.5), Inches(0.55),
              "P2: 6 areas   →   filter by tools   →   Area 3 + Area 6 kept   "
              "→   below: 4 active sub-directions in the literature",
              [],
              title_size=14, body_size=12)

acc_header = ["#", "Sub-direction", "From P2", "Representative Work", "Tools"]
acc = [
    ["1",
     [("Configuration studies", True)],
     [("Area 6", True)],
     "Schnepf 2023; Beier 2024 (semi-sub vs spar)",
     "OrcaFlex / OpenFAST"],
    ["2",
     [("Platform–cable global coupling", True)],
     [("Area 6", True)],
     "Sobhaniasl 2020; Hall 2021",
     "OrcaFlex / OpenFAST + MoorDyn"],
    ["3",
     [("Full fatigue workflow", True)],
     [("Area 3", True)],
     [("Beier 2023", True), (" (→ P5); Bakken 2019", False)],
     "OrcaFlex + UFLEX / Helica"],
    ["4",
     [("Configuration optimization", True)],
     [("Area 6", True)],
     "Rentschler 2019, 2020",
     "OrcaFlex + GA"],
]
add_table(s2, Inches(0.4), Inches(2.0),
          Inches(12.5), Inches(2.8),
          [acc_header] + acc,
          [Inches(0.4), Inches(3.2), Inches(1.2), Inches(4.4), Inches(3.3)],
          header_size=13, body_size=13)

add_textbox(s2, Inches(0.4), Inches(4.95), Inches(12.5), Inches(0.35),
            "Research gaps from §5 most relevant to us",
            size=14, bold=True, color=COLOR_TITLE)
gaps = [
    ("Gap A — ", "Linearization in global-local linking",
     ": σ = Kt·T + Kc·C ignores stick-slip nonlinearity."),
    ("Gap B — ", "Novel / deep-water configurations",
     ": fully suspended, W-shape, floater-type effects under-studied."),
    ("Gap C — ", "Project-specific designs",
     ": no rapid-screening tool for early-stage exploration."),
]
add_bullets(s2, Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.7),
            gaps, size=13)

add_footer(s2, 3)


# ============================================================
# P3: Our Feasible Path (was P2, simplified)
# ============================================================
s3 = prs.slides.add_slide(blank)

add_textbox(s3, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.55),
            "Our Feasible Path — OrcaFlex + AI",
            size=28, bold=True, color=COLOR_TITLE)
add_textbox(s3, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4),
            "Lab: OrcaFlex + Rhino + PC   |   Personal direction: AI algorithms",
            size=13, italic=True, color=COLOR_SOFT)

# Left: tool gate
add_textbox(s3, Inches(0.4), Inches(1.45), Inches(6.1), Inches(0.35),
            "① Tool-gate filter",
            size=15, bold=True, color=COLOR_TITLE)
gate_h = ["Direction", "We have", "We lack", "Feasible?"]
gate = [
    ["① Configuration", "OrcaFlex, Rhino", "—", "✅"],
    ["② Coupling", "OrcaFlex (6 DOF)", "—", "✅"],
    ["③ Fatigue", "OrcaFlex", "UFLEX",
     [("⚠️ → P4", True)]],
    ["④ Optimization", "OrcaFlex + Python", "—", "✅"],
]
add_table(s3, Inches(0.4), Inches(1.85),
          Inches(6.1), Inches(2.8),
          [gate_h] + gate,
          [Inches(1.6), Inches(1.8), Inches(1.4), Inches(1.3)],
          header_size=13, body_size=13)

# Right: AI candidates (shorter)
add_textbox(s3, Inches(6.8), Inches(1.45), Inches(6.1), Inches(0.35),
            "② AI entry points (question bank)",
            size=15, bold=True, color=COLOR_TITLE)
ai_h = ["Candidate", "Gap", "AI role"]
ai = [
    [[("A. Bending Kc correction", True)],
     "A", "Supervised / PINN"],
    [[("B. OrcaFlex surrogate", True)],
     "C", "NN surrogate"],
    [[("C. Configuration optimization", True)],
     "B", "Bayesian opt. / RL"],
    [[("D. Floater × sea-state sweep", True)],
     "B", "Active learning"],
]
add_table(s3, Inches(6.8), Inches(1.85),
          Inches(6.1), Inches(2.8),
          [ai_h] + ai,
          [Inches(2.8), Inches(1.0), Inches(2.3)],
          header_size=13, body_size=13)

# Bottom: single conclusion
add_highlight(s3, Inches(0.4), Inches(4.95),
              Inches(12.5), Inches(1.8),
              "Funnel conclusion",
              [
                  "Feasible research route = OrcaFlex global simulation + AI algorithms.",
                  "Entry into the literature: gaps A (linking), B (configurations), C (rapid screening).",
                  "Next slide: Beier 2023 — proof that fatigue analysis is possible without UFLEX.",
              ],
              title_size=15, body_size=13)

add_footer(s3, 4)


# ============================================================
# P4: Beier et al. (2023) — simplified
# ============================================================
s4 = prs.slides.add_slide(blank)

add_textbox(s4, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.55),
            "Beier et al. (2023) — Fatigue with Simplified Stress Factors",
            size=24, bold=True, color=COLOR_TITLE)
add_textbox(s4, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4),
            "Inter-array cable between two FOWTs  |  "
            "OC3-Hywind 5 MW  |  North Sea, 320 m  |  JMSE 11(7)",
            size=13, italic=True, color=COLOR_SOFT)

# Left: setup (condensed) + 4-step workflow
add_textbox(s4, Inches(0.4), Inches(1.4), Inches(6.1), Inches(0.35),
            "Setup",
            size=15, bold=True, color=COLOR_TITLE)
setup = [
    "Cable: 1260 m, fully suspended, 3 buoys @ 300 m",
    "Location: 61°N North Sea, 320 m depth",
    "Load cases: 30 sea states × 6 directions = 180 OrcaFlex runs",
    "Spectra: Torsethaugen (waves), NPD (wind)",
]
add_bullets(s4, Inches(0.4), Inches(1.8), Inches(6.1), Inches(1.7),
            setup, size=13)

add_textbox(s4, Inches(0.4), Inches(3.55), Inches(6.1), Inches(0.35),
            "4-step fatigue workflow",
            size=15, bold=True, color=COLOR_TITLE)
steps = [
    ("Step 1 — ", "UFLEX", ": nonlinear M(κ)"),
    ("Step 2 — ", "OrcaFlex", ": global T(t), C(t)"),
    ("Step 3 — ", "Stress factors Kt, Kc", ": σ(t) = Kt·T + Kc·C"),
    ("Step 4 — ", "Rainflow + Miner", ": fatigue damage"),
]
add_bullets(s4, Inches(0.4), Inches(3.95), Inches(6.1), Inches(2.0),
            steps, size=13)

# Right: innovation + findings
add_highlight(s4, Inches(6.8), Inches(1.4),
              Inches(6.1), Inches(1.85),
              "★ Key innovation",
              [
                  "Simplified analytical stress factors —",
                  "no UFLEX needed at preliminary design.",
                  "  • Kt (copper) = 232.3 kPa/kN — matches UFLEX exactly",
                  "  • Kc — up to 218% conservative for bending",
              ],
              title_size=14, body_size=13)

add_textbox(s4, Inches(6.8), Inches(3.4), Inches(6.1), Inches(0.35),
            "Key findings",
            size=15, bold=True, color=COLOR_TITLE)
findings = [
    [("Min fatigue life: ", False), ("≈ 7 × 10⁴ years", True),
     (" (≫ 25-year design life)", False)],
    [("Critical locations: ", False),
     ("hang-off + buoy attachment points", True), ("", False)],
    [("Failure mode: ", False), ("copper conductor first", True), ("", False)],
    [("Dominant driver: ", False), ("bending", True),
     (" (not axial tension)", False)],
    [("Marine growth (EOL): ", False), ("≈ 10% life reduction", True),
     ("", False)],
]
# Use bullets with rich runs
tb = s4.shapes.add_textbox(Inches(6.8), Inches(3.8),
                            Inches(6.1), Inches(2.2))
tf = tb.text_frame
tf.word_wrap = True
for i, parts in enumerate(findings):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(4)
    r0 = p.add_run()
    set_run(r0, "• ", size=13, color=COLOR_TEXT)
    for txt, b in parts:
        r = p.add_run()
        col = COLOR_ACCENT if b else COLOR_TEXT
        set_run(r, txt, size=13, bold=b, color=col)

add_highlight(s4, Inches(0.4), Inches(6.05),
              Inches(12.5), Inches(1.0),
              "Why this matters for us",
              [
                  "Beier proves a complete fatigue analysis is feasible with OrcaFlex alone. "
                  "The 218% bending conservatism is the exact opening for our AI work (→ P5).",
              ],
              title_size=14, body_size=13)

add_footer(s4, 5)


# ============================================================
# P5: Beyond Beier — simplified
# ============================================================
s5 = prs.slides.add_slide(blank)

add_textbox(s5, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.55),
            "Research Directions Beyond Beier (2023)",
            size=28, bold=True, color=COLOR_TITLE)
add_textbox(s5, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4),
            "Six extensions, all reachable with OrcaFlex + AI",
            size=13, italic=True, color=COLOR_SOFT)

beyond_h = ["#", "Direction", "Still open", "Our extension"]
beyond = [
    ["1", [("Bending Kc accuracy", True)],
     "True Kc still needs UFLEX",
     [("AI surrogate from literature data", True), (" — most direct", False)]],
    ["2", [("Floater type", True)],
     "TLP, barge untouched",
     "Multi-floater sampling"],
    ["3", [("Configuration variants", True)],
     "Lazy / W-shape; buoy layout",
     "Bayesian opt. / RL"],
    ["4", [("Environmental coverage", True)],
     "Other locations; misaligned wind/wave",
     "Active learning on hi-dim space"],
    ["5", [("Workflow acceleration", True)],
     "Each design = full re-sim",
     [("NN surrogate ", True), ("(~1000× faster)", False)]],
    ["6", [("Beyond fatigue", True)],
     "Installation, ULS, full coupling",
     "OrcaFlex direct; AI for UQ"],
]
add_table(s5, Inches(0.4), Inches(1.45),
          Inches(12.5), Inches(4.3),
          [beyond_h] + beyond,
          [Inches(0.5), Inches(2.7), Inches(4.2), Inches(5.1)],
          header_size=13, body_size=13)

add_highlight(s5, Inches(0.4), Inches(5.95),
              Inches(12.5), Inches(1.1),
              "★ Most promising next step",
              [
                  "Combine direction 1 + 5 — train an NN on Beier-style OrcaFlex outputs to (a) correct bending Kc",
                  "and (b) act as a surrogate for new designs. Discussion: which direction fits the group's research line?",
              ],
              title_size=14, body_size=13)

add_footer(s5, 6)


# ============================================================
# AI-P1 (overall P7): Code Flow & Key Technical Nodes (dedicated)
# ============================================================
s_flow = prs.slides.add_slide(blank)

# Section divider strip on top
fdiv = s_flow.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0),
                               SLIDE_W, Inches(0.15))
fdiv.fill.solid()
fdiv.fill.fore_color.rgb = COLOR_ACCENT
fdiv.line.fill.background()
fdiv.shadow.inherit = False

add_textbox(s_flow, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.4),
            "PART II — AI", size=13, bold=True, color=COLOR_ACCENT)
add_textbox(s_flow, Inches(0.4), Inches(0.6), Inches(12.5), Inches(0.55),
            "Mini Project — Purpose, Data & Pipeline",
            size=26, bold=True, color=COLOR_TITLE)
add_textbox(s_flow, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.35),
            "Why this project · what data we use · how the code flows",
            size=13, italic=True, color=COLOR_SOFT)

# ===== LEFT COLUMN: Purpose + Data Background =====
add_textbox(s_flow, Inches(0.4), Inches(1.6), Inches(5.8), Inches(0.35),
            "Purpose", size=15, bold=True, color=COLOR_TITLE)
purpose_items = [
    ("", "Step 0",
     " of the AI half of our cable-research funnel."),
    ("", "Build the pipeline",
     " (Pandas time-series → corr → curve) we'll reuse on OrcaFlex T(t), C(t) later."),
    ("", "Practice on public wind data first",
     "; cable surrogate (CNN-GRU) comes next."),
]
add_bullets(s_flow, Inches(0.4), Inches(1.95), Inches(5.8), Inches(1.8),
            purpose_items, size=12, space_after=5)

add_textbox(s_flow, Inches(0.4), Inches(3.85), Inches(5.8), Inches(0.35),
            "Dataset background", size=15, bold=True, color=COLOR_TITLE)
ds_h = ["Item", "Value"]
ds_rows = [
    ["Source", "Kaggle — Wind Turbine SCADA Dataset (public)"],
    ["Turbine", "Single onshore unit, Turkey"],
    ["Year", "2018, full year"],
    ["Sampling", "10-min interval"],
    ["Size", "50,460 rows × 5 columns"],
    ["Columns", "wind speed, power, theoretical, wind dir"],
    ["Quality", "0 % missing values"],
]
add_table(s_flow, Inches(0.4), Inches(4.2),
          Inches(5.8), Inches(2.6),
          [ds_h] + ds_rows,
          [Inches(1.35), Inches(4.45)],
          header_size=12, body_size=11)

# Right-column header
add_textbox(s_flow, Inches(6.6), Inches(1.6), Inches(6.4), Inches(0.35),
            "Code logic flowchart",
            size=15, bold=True, color=COLOR_TITLE)

# ---- Flowchart shape & arrow helpers (local to this slide) ----
FC_TERM_FILL = RGBColor(0xFF, 0xE0, 0xB2)
FC_PROC_FILL = RGBColor(0xE3, 0xF2, 0xFD)
FC_DEC_FILL = RGBColor(0xFF, 0xF1, 0xC5)
FC_IO_FILL = RGBColor(0xE8, 0xF5, 0xE9)
FC_ARROW = RGBColor(0x4A, 0x5A, 0x6E)

def flow_box(slide, shape_type, cx, cy, w, h, text, fill, text_size=11):
    left = Inches(cx - w / 2)
    top = Inches(cy - h / 2)
    shape = slide.shapes.add_shape(shape_type, left, top,
                                    Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = COLOR_TITLE
    shape.line.width = Pt(1.0)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(40000)
    tf.margin_right = Emu(40000)
    tf.margin_top = Emu(10000)
    tf.margin_bottom = Emu(10000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text, size=text_size, bold=True, color=COLOR_TEXT)
    return shape

def fc_v_arrow(slide, cx, top_y, length):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                    Inches(cx - 0.08), Inches(top_y),
                                    Inches(0.16), Inches(length))
    shape.fill.solid()
    shape.fill.fore_color.rgb = FC_ARROW
    shape.line.fill.background()
    shape.shadow.inherit = False

def fc_h_arrow(slide, left_x, cy, length, direction="right"):
    sh = (MSO_SHAPE.RIGHT_ARROW if direction == "right"
          else MSO_SHAPE.LEFT_ARROW)
    shape = slide.shapes.add_shape(sh,
                                    Inches(left_x), Inches(cy - 0.08),
                                    Inches(length), Inches(0.16))
    shape.fill.solid()
    shape.fill.fore_color.rgb = FC_ARROW
    shape.line.fill.background()
    shape.shadow.inherit = False

def fc_vline(slide, cx, top_y, bot_y):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(cx - 0.02), Inches(top_y),
                                    Inches(0.04), Inches(bot_y - top_y))
    shape.fill.solid()
    shape.fill.fore_color.rgb = FC_ARROW
    shape.line.fill.background()
    shape.shadow.inherit = False

def fc_hline(slide, left_x, right_x, cy):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left_x), Inches(cy - 0.02),
                                    Inches(right_x - left_x), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = FC_ARROW
    shape.line.fill.background()
    shape.shadow.inherit = False

# ---- Layout coordinates (right column only) ----
MAIN_X = 9.0
SIDE_X = 11.9
ARROW_LEN = 0.15

# (y_center, h, shape_type, fill, width, text)
nodes = [
    (2.10, 0.28, MSO_SHAPE.FLOWCHART_TERMINATOR, FC_TERM_FILL, 1.2, "Start"),
    (2.55, 0.32, MSO_SHAPE.FLOWCHART_DATA, FC_IO_FILL, 3.4,
     "Input · T1.csv (50,460 × 5)"),
    (2.99, 0.32, MSO_SHAPE.FLOWCHART_PROCESS, FC_PROC_FILL, 3.4,
     "Read CSV → DatetimeIndex"),
    (3.56, 0.50, MSO_SHAPE.FLOWCHART_DECISION, FC_DEC_FILL, 2.4,
     "missing > 5%?"),
    (4.32, 0.32, MSO_SHAPE.FLOWCHART_PROCESS, FC_PROC_FILL, 3.4,
     "Resample 10-min → 1-day"),
    (4.79, 0.32, MSO_SHAPE.FLOWCHART_PROCESS, FC_PROC_FILL, 3.4,
     "24-h rolling (center=True)"),
    (5.26, 0.32, MSO_SHAPE.FLOWCHART_PROCESS, FC_PROC_FILL, 3.4,
     "corr() + power-curve scatter"),
    (5.73, 0.32, MSO_SHAPE.FLOWCHART_DATA, FC_IO_FILL, 3.4,
     "Output · PNG @ 150 DPI"),
    (6.18, 0.28, MSO_SHAPE.FLOWCHART_TERMINATOR, FC_TERM_FILL, 1.2, "End"),
]

# Draw nodes
for cy, h, st, fill, w, text in nodes:
    flow_box(s_flow, st, MAIN_X, cy, w, h, text, fill,
             text_size=11)

# Vertical arrows between main nodes (skip decision→resample, handled below)
for i in range(len(nodes) - 1):
    if i == 3:
        continue
    cy_a, h_a, *_ = nodes[i]
    cy_b, h_b, *_ = nodes[i + 1]
    top = cy_a + h_a / 2
    bot = cy_b - h_b / 2
    if bot - top > 0.01:
        fc_v_arrow(s_flow, MAIN_X, top, bot - top)

# Side branch (yes path)
decision_cy = nodes[3][0]
decision_h = nodes[3][1]
decision_right = MAIN_X + nodes[3][4] / 2  # decision right edge
side_w = 1.6
side_left = SIDE_X - side_w / 2

flow_box(s_flow, MSO_SHAPE.FLOWCHART_PROCESS, SIDE_X, decision_cy,
         side_w, 0.32, "ffill / drop", FC_PROC_FILL, text_size=10)

# Yes arrow: decision → side box
fc_h_arrow(s_flow, decision_right, decision_cy,
           side_left - decision_right, direction="right")
add_textbox(s_flow, Inches(decision_right + 0.1),
            Inches(decision_cy - 0.28),
            Inches(0.6), Inches(0.22),
            "yes", size=10, bold=True, color=COLOR_ACCENT)

# "no" label below decision
decision_bot = decision_cy + decision_h / 2
add_textbox(s_flow, Inches(MAIN_X + 0.10), Inches(decision_bot + 0.0),
            Inches(0.4), Inches(0.2),
            "no", size=10, bold=True, color=COLOR_ACCENT)

# Decision → resample merge (combined no + side-return paths)
resample_top = nodes[4][0] - nodes[4][1] / 2
merge_y = (decision_bot + resample_top) / 2

# No path: vertical down from decision_bot to merge_y
fc_vline(s_flow, MAIN_X, decision_bot, merge_y)
# Yes return: vertical down from side box bottom to merge_y
side_bot = decision_cy + 0.32 / 2
fc_vline(s_flow, SIDE_X, side_bot, merge_y)
# Horizontal connect at merge_y
fc_hline(s_flow, MAIN_X, SIDE_X, merge_y)
# Final down arrow into resample top
fc_v_arrow(s_flow, MAIN_X, merge_y, resample_top - merge_y)

# ---- Legend strip (right column only) ----
legend_y = 6.55
legend_box = s_flow.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(6.6), Inches(legend_y),
                                      Inches(6.4), Inches(0.22))
legend_box.fill.solid()
legend_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
legend_box.line.color.rgb = COLOR_SOFT
legend_box.line.width = Pt(0.5)
legend_box.shadow.inherit = False
ltf = legend_box.text_frame
ltf.margin_left = Emu(100000)
ltf.margin_right = Emu(100000)
ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
lp = ltf.paragraphs[0]
lp.alignment = PP_ALIGN.CENTER
items = [
    ("Legend  ", True, COLOR_ACCENT),
    ("●", False, RGBColor(0xFF, 0xC0, 0x80)),
    (" Start / End      ", False, COLOR_TEXT),
    ("■", False, RGBColor(0x90, 0xC8, 0xE8)),
    (" Process      ", False, COLOR_TEXT),
    ("◆", False, RGBColor(0xE8, 0xD0, 0x80)),
    (" Decision      ", False, COLOR_TEXT),
    ("▱", False, RGBColor(0xA0, 0xD0, 0xA0)),
    (" Input / Output", False, COLOR_TEXT),
]
for txt, b, c in items:
    r = lp.add_run()
    set_run(r, txt, size=11, bold=b, color=c)

add_footer(s_flow, 7, total=9, section="AI Part")


# ============================================================
# AI-P2 (overall P8): Mini Project + EDA Findings
# ============================================================
s6 = prs.slides.add_slide(blank)

# Section divider strip on top
divider = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(0),
                              SLIDE_W, Inches(0.15))
divider.fill.solid()
divider.fill.fore_color.rgb = COLOR_ACCENT
divider.line.fill.background()
divider.shadow.inherit = False

add_textbox(s6, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.4),
            "PART II — AI", size=13, bold=True, color=COLOR_ACCENT)
add_textbox(s6, Inches(0.4), Inches(0.6), Inches(12.5), Inches(0.55),
            "EDA Findings (1/2) — Time-Series Views",
            size=26, bold=True, color=COLOR_TITLE)
add_textbox(s6, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.35),
            "Plot 1 daily mean wind  ·  Plot 2 24-h rolling smoothing",
            size=13, italic=True, color=COLOR_SOFT)

# ----- LEFT: Plot 1 -----
p1_path = PNG_DIR / "plot_1_daily_wind.png"
if p1_path.exists():
    s6.shapes.add_picture(str(p1_path), Inches(0.4), Inches(1.6),
                          width=Inches(6.1), height=Inches(3.05))
else:
    add_textbox(s6, Inches(0.4), Inches(3.0), Inches(6.1), Inches(0.5),
                "[plot_1 placeholder]", size=12, italic=True,
                color=COLOR_SOFT, align=PP_ALIGN.CENTER)

# Plot 1 caption + brief reading
add_textbox(s6, Inches(0.4), Inches(4.75), Inches(6.1), Inches(0.35),
            "Plot 1 · Daily mean wind speed",
            size=14, bold=True, color=COLOR_TITLE)
p1_lines = [
    ("", "Year mean = 7.56 m/s",
     " (dashed red line on plot)."),
    ("", "Winter peak ~10 m/s",
     " (Nov-Feb); summer trough ~4 m/s (Jun-Aug)."),
    ("", "Take-away:",
     " annual output strongly winter-skewed."),
]
add_bullets(s6, Inches(0.4), Inches(5.10), Inches(6.1), Inches(1.7),
            p1_lines, size=12, space_after=4)

# ----- RIGHT: Plot 2 -----
p2_path = PNG_DIR / "plot_2_rolling.png"
if p2_path.exists():
    s6.shapes.add_picture(str(p2_path), Inches(6.85), Inches(1.6),
                          width=Inches(6.1), height=Inches(3.05))
else:
    add_textbox(s6, Inches(6.85), Inches(3.0), Inches(6.1), Inches(0.5),
                "[plot_2 placeholder]", size=12, italic=True,
                color=COLOR_SOFT, align=PP_ALIGN.CENTER)

# Plot 2 caption + brief reading
add_textbox(s6, Inches(6.85), Inches(4.75), Inches(6.1), Inches(0.35),
            "Plot 2 · 24-h rolling mean",
            size=14, bold=True, color=COLOR_TITLE)
p2_lines = [
    ("", "window = 24 (one diurnal cycle)",
     ", center=True (no phase lag)."),
    ("", "Hourly noise removed",
     "; seasonal trend visible (red line)."),
    ("", "Take-away:",
     " use this curve to flag sensor drift or long calm spells."),
]
add_bullets(s6, Inches(6.85), Inches(5.10), Inches(6.1), Inches(1.7),
            p2_lines, size=12, space_after=4)

# Bottom note
add_highlight(s6, Inches(0.4), Inches(6.85),
              Inches(12.5), Inches(0.25),
              "★ Next slide: Plot 3 (correlation) + Plot 4 (power curve).",
              [],
              title_size=12, body_size=11)

add_footer(s6, 8, total=9, section="AI Part")


# ============================================================
# AI-P2 (overall P8): 12-Week Roadmap
# ============================================================
s8 = prs.slides.add_slide(blank)

# Section strip
divider3 = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0),
                               SLIDE_W, Inches(0.15))
divider3.fill.solid()
divider3.fill.fore_color.rgb = COLOR_ACCENT
divider3.line.fill.background()
divider3.shadow.inherit = False

add_textbox(s8, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.4),
            "PART II — AI", size=13, bold=True, color=COLOR_ACCENT)
add_textbox(s8, Inches(0.4), Inches(0.6), Inches(12.5), Inches(0.55),
            "EDA Findings (2/2) — Statistical & Curve Views",
            size=26, bold=True, color=COLOR_TITLE)
add_textbox(s8, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.35),
            "Plot 3 correlation matrix  ·  Plot 4 power curve",
            size=13, italic=True, color=COLOR_SOFT)

# ----- LEFT: Plot 3 (correlation) -----
p3_path = PNG_DIR / "plot_3_corr.png"
if p3_path.exists():
    s8.shapes.add_picture(str(p3_path), Inches(0.6), Inches(1.55),
                          width=Inches(5.4), height=Inches(3.5))
else:
    add_textbox(s8, Inches(0.4), Inches(3.0), Inches(6.0), Inches(0.5),
                "[plot_3 placeholder]", size=12, italic=True,
                color=COLOR_SOFT, align=PP_ALIGN.CENTER)

add_textbox(s8, Inches(0.4), Inches(5.15), Inches(6.1), Inches(0.35),
            "Plot 3 · Correlation matrix (Pearson)",
            size=14, bold=True, color=COLOR_TITLE)
p3_lines = [
    ("", "wind_speed ↔ power_kW = 0.91",
     " (boxed cell on plot)."),
    ("", "wind_dir correlations ≈ 0",
     " — yaw control makes direction irrelevant."),
    ("", "Take-away:",
     " wind speed is THE dominant predictor."),
]
add_bullets(s8, Inches(0.4), Inches(5.50), Inches(6.1), Inches(1.4),
            p3_lines, size=12, space_after=4)

# ----- RIGHT: Plot 4 (power curve) -----
p4_path = PNG_DIR / "plot_4_power_curve.png"
if p4_path.exists():
    s8.shapes.add_picture(str(p4_path), Inches(6.85), Inches(1.55),
                          width=Inches(6.1), height=Inches(3.5))
else:
    add_textbox(s8, Inches(6.85), Inches(3.0), Inches(6.1), Inches(0.5),
                "[plot_4 placeholder]", size=12, italic=True,
                color=COLOR_SOFT, align=PP_ALIGN.CENTER)

add_textbox(s8, Inches(6.85), Inches(5.15), Inches(6.1), Inches(0.35),
            "Plot 4 · Power curve (measured vs theoretical)",
            size=14, bold=True, color=COLOR_TITLE)
p4_lines = [
    ("", "Cubic rise 3-12 m/s",
     "; saturation at rated 3600 kW."),
    ("", "Red-shaded zone above 12 m/s",
     ": many points below theoretical → losses."),
    ("", "Take-away:",
     " losses zone = curtailment / faults, future modelling target."),
]
add_bullets(s8, Inches(6.85), Inches(5.50), Inches(6.1), Inches(1.4),
            p4_lines, size=12, space_after=4)

# Bottom note
add_highlight(s8, Inches(0.4), Inches(6.85),
              Inches(12.5), Inches(0.25),
              "★ Tonight's deliverable: a clean, narrated EDA — same "
              "template will be reused on OrcaFlex T(t), C(t) outputs.",
              [],
              title_size=12, body_size=11)

add_footer(s8, 9, total=9, section="AI Part")


# ---------- Save ----------
prs.save(SLIDES_OUT)
print(f"Saved: {SLIDES_OUT}")
print(f"Slides: {len(prs.slides)}")
